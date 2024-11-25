from metadata import *
import PyPDF2
import docx
from PIL import Image
import io
import fitz
from transformers import BlipProcessor, BlipForConditionalGeneration
import pptx
import moviepy.editor as mp
from moviepy.video.io.ffmpeg_tools import ffmpeg_extract_subclip
import whisper
import os

processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
model = BlipForConditionalGeneration.from_pretrained(
    "Salesforce/blip-image-captioning-base"
)


def generate_caption(image: Image.Image) -> str:
    """Generate a caption for an image using a pre-trained model."""
    inputs = processor(image, return_tensors="pt")
    output = model.generate(**inputs)
    caption = processor.decode(output[0], skip_special_tokens=True)
    return caption


def extract_word(file):
    doc = docx.Document(file)
    text = "\n".join([para.text for para in doc.paragraphs])
    image_captions = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            image = Image.open(io.BytesIO(image_data))

            caption = generate_caption(image)
            image_captions.append(caption)

    if len(image_captions) > 1:
        image_captions = image_captions[::-1]
    combined_output = text
    for caption in image_captions:
        combined_output += f"\nImage caption: {caption}"

    return combined_output


def extract_pdf(file):
    text_output = ""
    reader = PyPDF2.PdfReader(file)
    file.seek(0)
    pdf_document = fitz.open(stream=file.read(), filetype="pdf")

    for page_index in range(len(pdf_document)):
        text_output += reader.pages[page_index].extract_text() + "\n"

        page = pdf_document.load_page(page_index)
        image_list = page.get_images(full=True)

        for _, img in enumerate(image_list, start=1):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]

            image = Image.open(io.BytesIO(image_bytes))
            caption = generate_caption(image)
            text_output += "Image caption: " + caption + "\n"
    return text_output.strip()


def extract_pptx(file):
    # Open the PowerPoint file
    presentation = pptx.Presentation(file)

    # Initialize text and image outputs
    combined_output = ""

    # Process each slide
    for slide_index, slide in enumerate(presentation.slides):
        # Extract text from the slide
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text.append(paragraph.text)

        # Add slide text to the output
        if slide_text:
            combined_output += (
                f"Slide {slide_index + 1} Text:\n" + "\n".join(slide_text) + "\n"
            )

        # Extract images from the slide
        for shape in slide.shapes:
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                pil_image = Image.open(io.BytesIO(image_bytes))

                # Generate a caption for the image
                caption = generate_caption(pil_image)

                # Add the caption to the output
                combined_output += f"Slide {slide_index + 1} Image Caption: {caption}\n"

    return combined_output.strip()


whisper_model = whisper.load_model("base")


def extract_audio_with_timestamps(path):
    """Extracts and transcribes audio with timestamps from a video file."""
    # Extract audio from video
    video = mp.VideoFileClip(path)
    audio = video.audio
    script_dir = os.path.dirname(os.path.abspath(__file__))
    audio_path = os.path.join(script_dir, "temp_audio.wav")
    audio.write_audiofile(audio_path)
    audio.close()
    result = whisper_model.transcribe(audio_path, task="transcribe", verbose=False)
    segments = result["segments"]

    audio_transcription = []
    for segment in segments:
        start = segment["start"]
        end = segment["end"]
        text = segment["text"]
        audio_transcription.append((start, end, text.strip()))

    if os.path.exists(audio_path):
        os.remove(audio_path)
    return audio_transcription


def extract_video(file):
    """Extracts and combines video frame captions and audio transcription."""
    # Save video file locally
    script_dir = os.path.dirname(os.path.abspath(__file__))
    video_path = os.path.join(script_dir, "temp_video.mp4")
    with open(video_path, "wb") as f:
        f.write(file.read())

    try:
        audio_transcription = extract_audio_with_timestamps(video_path)

        video = mp.VideoFileClip(video_path)
        duration = int(video.duration)
        frame_interval = 5
        frame_captions = []
        # get video image every 5 seconds
        for t in range(0, duration, frame_interval):
            frame = video.get_frame(t)
            pil_image = Image.fromarray(frame)
            caption = generate_caption(pil_image)
            frame_captions.append((t, caption))

        # sort timelines so that audio and video aligns well
        timeline = []
        for start, end, text in audio_transcription:
            timeline.append((start, f"[Audio {start:.2f}-{end:.2f}s]: {text}"))

        for time, caption in frame_captions:
            timeline.append((time, f"[Video Frame {time}s]: {caption}"))

        timeline.sort(key=lambda x: x[0])

        text_output = "\n".join(event[1] for event in timeline)
    finally:
        if "video" in locals():
            video.close()
        if os.path.exists(video_path):
            os.remove(video_path)
        return text_output.strip()
